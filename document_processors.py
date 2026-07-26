# document_processors.py
#
# Extracts text/table/image/audio content from uploaded files and returns
# plain Python dicts ready for chunking + indexing. Everything happens
# in-memory (BytesIO / bytes) — nothing is ever written to the serverless
# function's read-only filesystem.

import io
import os
from typing import List, Tuple, Dict

import fitz  # PyMuPDF
from pptx import Presentation

from utils import (
    describe_image,
    is_graph,
    describe_chart_or_table,
    transcribe_audio,
    extract_text_around_item,
    process_text_blocks,
)

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".gif"}
AUDIO_EXTENSIONS = {".mp3", ".wav", ".m4a", ".flac", ".ogg", ".webm"}


def process_pdf(file_bytes: bytes, filename: str) -> List[Dict]:
    """Extract text passages, tables, and figures from a PDF, entirely
    in-memory."""
    docs: List[Dict] = []
    base_name = os.path.splitext(filename)[0]

    try:
        pdf = fitz.open(stream=file_bytes, filetype="pdf")
    except Exception as e:
        return [{
            "text": f"[Could not open PDF '{filename}': {e}]",
            "metadata": {"source": filename, "type": "error", "page_num": 0},
        }]

    for page_num in range(len(pdf)):
        page = pdf[page_num]

        text_blocks = [
            b for b in page.get_text("blocks", sort=True)
            if b[-1] == 0
            and not (b[1] < page.rect.height * 0.1 or b[3] > page.rect.height * 0.9)
        ]
        grouped_text_blocks = process_text_blocks(text_blocks)

        # --- Tables -----------------------------------------------------
        table_bboxes = []
        try:
            tables = page.find_tables(
                horizontal_strategy="lines_strict", vertical_strategy="lines_strict"
            )
            for idx, tab in enumerate(tables, start=1):
                try:
                    table_df = tab.to_pandas()
                    bbox = fitz.Rect(tab.bbox)
                    table_bboxes.append(bbox)

                    before_text, after_text = extract_text_around_item(
                        text_blocks, bbox, page.rect.height
                    )
                    columns = ", ".join(str(c) for c in table_df.columns)
                    caption = (before_text + " " + after_text).strip()
                    if not caption:
                        caption = columns

                    try:
                        table_repr = table_df.to_markdown(index=False)
                    except Exception:
                        table_repr = table_df.to_csv(index=False)

                    docs.append({
                        "text": (
                            f"This is a table with caption: {caption}\n"
                            f"Columns: {columns}\n{table_repr}"
                        ),
                        "metadata": {
                            "source": f"{base_name}-page{page_num}-table{idx}",
                            "type": "table",
                            "page_num": page_num,
                        },
                    })
                except Exception:
                    continue
        except Exception:
            pass

        # --- Images / charts ---------------------------------------------
        try:
            for image_info in page.get_image_info(xrefs=True):
                xref = image_info.get("xref", 0)
                if xref == 0:
                    continue

                img_bbox = fitz.Rect(image_info["bbox"])
                if (
                    img_bbox.width < page.rect.width / 20
                    or img_bbox.height < page.rect.height / 20
                ):
                    continue

                extracted = pdf.extract_image(xref)
                image_bytes = extracted["image"]

                before_text, after_text = extract_text_around_item(
                    text_blocks, img_bbox, page.rect.height
                )
                caption_hint = (before_text + " " + after_text).strip()

                try:
                    description = describe_image(image_bytes)
                    if is_graph(description):
                        description = describe_chart_or_table(image_bytes)
                except Exception as e:
                    description = f"[Image description unavailable: {e}]"

                docs.append({
                    "text": (
                        f"This is an image on page {page_num}"
                        + (f" (context: {caption_hint})" if caption_hint else "")
                        + f". Description: {description}"
                    ),
                    "metadata": {
                        "source": f"{base_name}-page{page_num}-image{xref}",
                        "type": "image",
                        "page_num": page_num,
                    },
                })
        except Exception:
            pass

        # --- Body text ------------------------------------------------------
        for block_ctr, (heading_block, content) in enumerate(grouped_text_blocks, start=1):
            heading_bbox = fitz.Rect(heading_block[:4])
            if any(heading_bbox.intersects(tb) for tb in table_bboxes):
                continue
            docs.append({
                "text": f"{heading_block[4]}\n{content}",
                "metadata": {
                    "source": f"{base_name}-page{page_num}-block{block_ctr}",
                    "type": "text",
                    "page_num": page_num,
                },
            })

    pdf.close()
    return docs


def process_image_file(file_bytes: bytes, filename: str) -> List[Dict]:
    try:
        description = describe_image(file_bytes)
    except Exception as e:
        description = f"[Image description unavailable: {e}]"
    return [{
        "text": description,
        "metadata": {"source": filename, "type": "image", "page_num": 0},
    }]


def process_pptx_file(file_bytes: bytes, filename: str) -> List[Dict]:
    """Extract slide text + speaker notes from a .pptx file.

    Note: legacy binary .ppt and slide-to-image rendering (which the
    original app did via a local LibreOffice install) are not supported
    here, since Vercel's serverless functions have no LibreOffice binary
    available. Ask users to export legacy .ppt files to .pptx first.
    """
    docs: List[Dict] = []
    try:
        prs = Presentation(io.BytesIO(file_bytes))
    except Exception as e:
        return [{
            "text": f"[Could not open presentation '{filename}': {e}]",
            "metadata": {"source": filename, "type": "error", "page_num": 0},
        }]

    for slide_num, slide in enumerate(prs.slides):
        texts = [
            shape.text for shape in slide.shapes
            if hasattr(shape, "text") and shape.text
        ]
        slide_text = "\n".join(texts)
        try:
            notes = (
                slide.notes_slide.notes_text_frame.text
                if slide.has_notes_slide else ""
            )
        except Exception:
            notes = ""

        combined = slide_text + (f"\n\nSpeaker notes: {notes}" if notes else "")
        if combined.strip():
            docs.append({
                "text": combined,
                "metadata": {
                    "source": f"{filename}-slide{slide_num + 1}",
                    "type": "slide",
                    "page_num": slide_num,
                },
            })
    return docs


def process_text_file(file_bytes: bytes, filename: str) -> List[Dict]:
    try:
        text = file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        text = file_bytes.decode("latin-1", errors="ignore")
    return [{
        "text": text,
        "metadata": {"source": filename, "type": "text", "page_num": 0},
    }]


def process_audio_file(file_bytes: bytes, filename: str) -> List[Dict]:
    try:
        transcript = transcribe_audio(file_bytes, filename)
    except Exception as e:
        transcript = f"[Audio transcription unavailable: {e}]"
    return [{
        "text": transcript,
        "metadata": {"source": filename, "type": "audio", "page_num": 0},
    }]


def load_multimodal_data(files: List[Tuple[str, bytes]]) -> List[Dict]:
    """Route each uploaded (filename, bytes) pair to the right processor.

    Returns a flat list of {"text": ..., "metadata": {...}} documents, ready
    to be chunked and indexed.
    """
    documents: List[Dict] = []
    for filename, file_bytes in files:
        ext = os.path.splitext(filename.lower())[1]
        try:
            if ext in IMAGE_EXTENSIONS:
                documents.extend(process_image_file(file_bytes, filename))
            elif ext == ".pdf":
                documents.extend(process_pdf(file_bytes, filename))
            elif ext == ".pptx":
                documents.extend(process_pptx_file(file_bytes, filename))
            elif ext in AUDIO_EXTENSIONS:
                documents.extend(process_audio_file(file_bytes, filename))
            else:
                documents.extend(process_text_file(file_bytes, filename))
        except Exception as e:
            documents.append({
                "text": f"[Error processing {filename}: {e}]",
                "metadata": {"source": filename, "type": "error", "page_num": 0},
            })
    return documents
